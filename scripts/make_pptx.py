"""
몽글픽 중간보고 PPTX 생성 스크립트

작업 1: 템플릿 기반 버전 (프로젝트_ppt_가이드.pptx 텍스트 교체)
작업 2: 신규 디자인 버전 (python-pptx로 처음부터 작성)

실행:
  PYTHONPATH=src uv run python scripts/make_pptx.py
"""

import shutil
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.enum.dml import MSO_THEME_COLOR

try:
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
except ImportError:
    pass

# ─── 경로 설정 ───────────────────────────────────────────────────────────────
DOCS_DIR = Path("/Users/yoonhyungjoo/Documents/monglepick/docs")
TEMPLATE_SRC = DOCS_DIR / "프로젝트_ppt_가이드.pptx"
TEMPLATE_DST = DOCS_DIR / "중간보고_몽글픽_템플릿버전.pptx"
NEW_DST = DOCS_DIR / "중간보고_몽글픽_신규버전.pptx"


# ─── 공통 유틸 ───────────────────────────────────────────────────────────────

def rgb_to_hex(color) -> str:
    """
    RGBColor(bytes 서브클래스) 또는 (r,g,b) 튜플을 6자리 hex 문자열로 변환한다.
    RGBColor는 bytes이므로 인덱스([0][1][2])로 각 채널에 접근한다.
    """
    r, g, b = color[0], color[1], color[2]
    return f"{r:02x}{g:02x}{b:02x}"

def set_text_keep_format(text_frame, new_text: str):
    """
    텍스트프레임의 첫 번째 paragraph/run 서식을 유지하면서
    전체 텍스트를 교체한다.
    줄바꿈은 \n 으로 표현하며 paragraph를 새로 추가한다.
    """
    # 기존 첫 번째 run 서식 백업
    first_para = text_frame.paragraphs[0]
    first_run_font = None
    if first_para.runs:
        r = first_para.runs[0]
        # 색상 타입이 RGB일 때만 추출 (SchemeColor 등은 None 처리)
        try:
            from pptx.enum.dml import MSO_THEME_COLOR
            color_val = r.font.color.rgb if (r.font.color and r.font.color.type and
                                              str(r.font.color.type) == "RGB (1)") else None
        except (AttributeError, Exception):
            color_val = None
        first_run_font = {
            "bold": r.font.bold,
            "size": r.font.size,
            "color": color_val,
            "name": r.font.name,
        }

    # 기존 paragraph 요소 모두 제거
    tf_elem = text_frame._txBody
    for p_elem in tf_elem.findall(qn("a:p")):
        tf_elem.remove(p_elem)

    lines = new_text.split("\n")
    for line_idx, line in enumerate(lines):
        from pptx.oxml import parse_xml
        from lxml import etree

        p_elem = etree.SubElement(tf_elem, qn("a:p"))
        r_elem = etree.SubElement(p_elem, qn("a:r"))

        # rPr (run properties) 설정
        rpr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ko-KR", "dirty": "0"})
        if first_run_font:
            if first_run_font.get("bold"):
                rpr.set("b", "1")
            if first_run_font.get("size"):
                rpr.set("sz", str(first_run_font["size"] // 100))
            if first_run_font.get("name"):
                latin = etree.SubElement(rpr, qn("a:latin"))
                latin.set("typeface", first_run_font["name"])
            if first_run_font.get("color"):
                solid_fill = etree.SubElement(rpr, qn("a:solidFill"))
                srg = etree.SubElement(solid_fill, qn("a:srgbClr"))
                srg.set("val", rgb_to_hex(first_run_font["color"]))

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line


def replace_text_in_shape(shape, old_text: str, new_text: str):
    """도형 또는 그룹 내에서 재귀적으로 텍스트를 교체한다."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            replace_text_in_shape(s, old_text, new_text)
    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
        current = shape.text_frame.text
        if old_text in current:
            set_text_keep_format(shape.text_frame, new_text)


def replace_text_direct(shape, new_text: str):
    """도형의 텍스트를 직접 교체 (정확히 이 shape를 대상으로)."""
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        set_text_keep_format(shape.text_frame, new_text)


def find_shape_by_name(slide, name: str):
    """슬라이드에서 이름으로 shape 검색 (그룹 내 재귀)."""
    def _find(shapes, name):
        for s in shapes:
            if s.name == name:
                return s
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                found = _find(s.shapes, name)
                if found:
                    return found
        return None
    return _find(slide.shapes, name)


def find_shape_by_id(slide, shape_id: int):
    """슬라이드에서 ID로 shape 검색 (그룹 내 재귀)."""
    def _find(shapes, sid):
        for s in shapes:
            if s.shape_id == sid:
                return s
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                found = _find(s.shapes, sid)
                if found:
                    return found
        return None
    return _find(slide.shapes, shape_id)


def get_table_shape(slide):
    """슬라이드에서 TABLE 타입 shape 반환."""
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            return shape
    return None


def set_cell_text(cell, text: str, font_size=None, bold=False, color=None, align=None):
    """테이블 셀 텍스트 설정 (서식 포함)."""
    tf = cell.text_frame
    tf.word_wrap = True
    from lxml import etree
    # 기존 내용 클리어
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    p_elem = etree.SubElement(txBody, qn("a:p"))

    # 정렬 설정
    if align:
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        align_map = {"center": "ctr", "left": "l", "right": "r"}
        pPr.set("algn", align_map.get(align, "l"))

    r_elem = etree.SubElement(p_elem, qn("a:r"))
    rpr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ko-KR", "dirty": "0"})
    if bold:
        rpr.set("b", "1")
    if font_size:
        rpr.set("sz", str(int(font_size * 100)))
    if color:
        solid_fill = etree.SubElement(rpr, qn("a:solidFill"))
        srg = etree.SubElement(solid_fill, qn("a:srgbClr"))
        srg.set("val", f"{color[0]:02x}{color[1]:02x}{color[2]:02x}")
    latin = etree.SubElement(rpr, qn("a:latin"))
    latin.set("typeface", "Malgun Gothic")

    t_elem = etree.SubElement(r_elem, qn("a:t"))
    t_elem.text = text


# ─── 작업 1: 템플릿 기반 버전 ────────────────────────────────────────────────

def make_template_version():
    """
    프로젝트_ppt_가이드.pptx를 복사하고 텍스트를 몽글픽 내용으로 교체한다.
    """
    print("[작업 1] 템플릿 기반 버전 생성 시작...")

    # 템플릿 복사
    shutil.copy(str(TEMPLATE_SRC), str(TEMPLATE_DST))
    prs = Presentation(str(TEMPLATE_DST))

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────────────
    slide1 = prs.slides[0]
    s = find_shape_by_id(slide1, 10)  # TextBox 9 "팀 프로젝트명(주제)"
    if s:
        replace_text_direct(s, "몽글픽 (MongLePick)\nAI 기반 개인화 영화 추천 서비스")

    # ── 슬라이드 3: 프로젝트 개요 ─────────────────────────────────────────────
    slide3 = prs.slides[2]

    # TextBox 77
    s = find_shape_by_id(slide3, 78)  # name='TextBox 77'
    if s:
        replace_text_direct(s, "영화 추천에서 출발한 AI 에이전트\n910,140건 DB + 하이브리드 RAG 검색\n5개 서비스 마이크로서비스 아키텍처")

    # 그룹 44 내 텍스트 교체
    g44 = find_shape_by_name(slide3, "그룹 44")
    if g44:
        for s in g44.shapes:
            if hasattr(s, "has_text_frame") and s.has_text_frame:
                t = s.text_frame.text
                if "프로젝트 주제 및 선정 배경" in t:
                    replace_text_direct(s, "AI 영화 추천 + LangGraph 에이전트\n사용자 취향 대화 기반 초개인화")
                elif "특화 포인트" in t:
                    replace_text_direct(s, "Qdrant+ES+Neo4j 하이브리드 RAG\n멀티턴 대화 + 이미지 분석 지원")

    # 그룹 90 내 텍스트 교체
    g90 = find_shape_by_name(slide3, "그룹 90")
    if g90:
        for s in g90.shapes:
            if hasattr(s, "has_text_frame") and s.has_text_frame:
                t = s.text_frame.text
                if "프로젝트\x0b내용" in t or "프로젝트\n내용" in t or t.strip() == "프로젝트 내용":
                    replace_text_direct(s, "서비스\n구현 내용")
                elif "구현 내용" in t or "컨셉" in t:
                    replace_text_direct(s, "추천/검색/결제/커뮤니티 통합")

    # 그룹 101 내 텍스트 교체
    g101 = find_shape_by_name(slide3, "그룹 101")
    if g101:
        for s in g101.shapes:
            if hasattr(s, "has_text_frame") and s.has_text_frame:
                t = s.text_frame.text
                if "활용 장비" in t:
                    replace_text_direct(s, "활용 기술\n스택")
                elif "개발환경" in t:
                    replace_text_direct(s, "Ollama, Qdrant, ES, Neo4j,\nMySQL, Redis, React")

    # 그룹 112 내 텍스트 교체
    g112 = find_shape_by_name(slide3, "그룹 112")
    if g112:
        for s in g112.shapes:
            if hasattr(s, "has_text_frame") and s.has_text_frame:
                t = s.text_frame.text
                if "활용방안" in t:
                    replace_text_direct(s, "활용방안 및\n기대 효과")
                elif "기대 효용" in t or "비즈니스" in t:
                    replace_text_direct(s, "영화 소비 패턴 개인화\n포인트 기반 AI 이용권 경제")

    # ── 슬라이드 4: 팀 구성 ───────────────────────────────────────────────────
    slide4 = prs.slides[3]

    # 표 수정
    tbl_shape = get_table_shape(slide4)
    if tbl_shape:
        tbl = tbl_shape.table
        # 헤더 행
        set_cell_text(tbl.cell(0, 0), "훈련생", bold=True, font_size=11)
        set_cell_text(tbl.cell(0, 1), "역할", bold=True, font_size=11)
        set_cell_text(tbl.cell(0, 2), "담당 업무", bold=True, font_size=11)
        # 데이터 행
        members = [
            ("윤형주", "팀장", "AI Agent(LangGraph), 추천엔진, 포인트/결제, 관리자, Client"),
            ("김민규", "팀원", "인증(JWT/OAuth2), 사용자, 플레이리스트, 대시보드"),
            ("이민수", "팀원", "커뮤니티(게시판/리뷰), 콘텐츠 관리"),
            ("정한나", "팀원", "추천서버(FastAPI), 통계/분석, 데이터 파이프라인"),
        ]
        for row_idx, (name, role, task) in enumerate(members, start=1):
            set_cell_text(tbl.cell(row_idx, 0), name, font_size=10)
            set_cell_text(tbl.cell(row_idx, 1), role, font_size=10)
            set_cell_text(tbl.cell(row_idx, 2), task, font_size=10)

    # TextBox 159 교체
    s = find_shape_by_id(slide4, 160)  # name='TextBox 159'
    if s:
        replace_text_direct(s, "React Client 구현")

    # ── 슬라이드 5: WBS ────────────────────────────────────────────────────────
    slide5 = prs.slides[4]

    # TextBox 4
    s = find_shape_by_id(slide5, 5)  # name='TextBox 4'
    if s:
        replace_text_direct(s, "6단계 Phase 기반 점진적 구현\n설계→데이터→AI→서비스→통합→고도화")

    # 표 수정
    tbl_shape = get_table_shape(slide5)
    if tbl_shape:
        tbl = tbl_shape.table
        wbs_data = [
            ("구분", "기간", "활동", "비고"),
            ("요구사항 분석", "3/3 ~ 3/7", "기획/DB설계/API설계", "팀 협의"),
            ("데이터 수집/적재", "3/10 ~ 3/14", "TMDB/KOBIS 수집, 5DB 적재", "91만건"),
            ("핵심 기능 구현", "3/17 ~ 3/28", "AI Agent, Backend, Client", "Phase 0~5"),
            ("통합/고도화", "3/31 ~ 4/4", "추천 품질, 관리자, 다크모드", "ML Phase 완료"),
            ("중간보고", "4/7", "중간 발표", "오늘"),
            ("최종 마무리", "4/8 ~ 4/18", "LoRA 파인튜닝, DB 재적재", "예정"),
        ]
        for row_idx, row_data in enumerate(wbs_data):
            for col_idx, cell_text in enumerate(row_data):
                is_header = row_idx == 0
                set_cell_text(
                    tbl.cell(row_idx, col_idx),
                    cell_text,
                    bold=is_header,
                    font_size=10 if not is_header else 11,
                )

    # WBS 그룹 내 텍스트 노드 교체
    wbs_group_replacements = {
        "프로젝트 기획 및 주제 선정": "요구사항 분석",
        "기획안 작성": "DB/API 설계",
        "필요 데이터  및 수집 절차 정의": "데이터 파이프라인",
        "필요 데이터 및 수집 절차 정의": "데이터 파이프라인",
        "외부 데이터 수집": "TMDB/KOBIS/KMDb",
        "데이터 정제 및 정규화": "벡터 임베딩/적재",
        "모형 구현": "AI Agent 구현",
        "모바일 서비스 시스템 설계": "Backend API 구현",
        "모바일 플랫폼 구현": "Client/Admin 구현",
    }
    for old, new in wbs_group_replacements.items():
        replace_text_in_shape_on_slide(slide5, old, new)

    # ── 슬라이드 6: 수행경과 ──────────────────────────────────────────────────
    slide6 = prs.slides[5]

    progress_replacements = {
        "결과를 서술하는 과정에서는 활용된 기술 및 구현 방법, 핵심기능, 구현 결과": "AI Agent 14노드 LangGraph 구현\n이미지 분석, 멀티턴 대화, 7개 Tools",
        "프로젝트 수행 과정이 잘 드러날 수 있도록 가공 과정부터 활용까지 전체적인 프로세스를 단계별로 작성": "910,140건 영화 DB + 하이브리드 RAG\nQdrant+ES+Neo4j RRF 검색 파이프라인",
        "프로젝트 수행 과정에서의 피드백 내용과 그것을 적용(보완 등)한 내용이 포함되도록 작성": "의도 분류 → 선호 추출 → 추천 생성\n포인트/결제/등급 시스템 통합",
        "프로젝트 수행 결과물을 잘 드러낼 수 있는 자료를 첨부하여 작성": "250+ Client 모듈, 749 Admin 모듈\n다크모드/반응형 완성",
        "결과물 사진, 시연 영상, 구동 화면 등 프로젝트의 우수성이 드러날 수 있는 자료": "",
    }
    for old, new in progress_replacements.items():
        replace_text_in_shape_on_slide(slide6, old, new)

    # ── 슬라이드 7: 결과 제시 ① ──────────────────────────────────────────────
    slide7 = prs.slides[6]
    s = find_shape_by_id(slide7, 16)  # TextBox 15
    if s:
        replace_text_direct(s, "* 결과 제시 ① 데이터 파이프라인")
    s = find_shape_by_id(slide7, 55)  # TextBox 9
    if s:
        replace_text_direct(s,
            "데이터 수집: TMDB API / Kaggle / KOBIS / KMDb\n"
            "총 910,140건 → 5개 DB 적재\n"
            "Upstage Solar 임베딩 4096차원\n"
            "Elasticsearch Nori 한국어 형태소 분석"
        )
    s = find_shape_by_id(slide7, 56)  # TextBox 19
    if s:
        replace_text_direct(s,
            "5DB 적재 현황\n"
            "- Qdrant: 벡터 인덱스 (4096차원, Cosine)\n"
            "- ES: BM25 역인덱스 (Nori 분석기)\n"
            "- Neo4j: 관계 그래프 (감독/배우/장르)\n"
            "- MySQL: 메타데이터 (58개 테이블)\n"
            "- Redis: 세션 + CF 캐시"
        )

    # ── 슬라이드 8: 결과 제시 ② ──────────────────────────────────────────────
    slide8 = prs.slides[7]
    s = find_shape_by_id(slide8, 16)  # TextBox 15
    if s:
        replace_text_direct(s, "* 결과 제시 ② AI Agent 아키텍처")
    s = find_shape_by_id(slide8, 25)  # TextBox 1
    if s:
        replace_text_direct(s,
            "LangGraph StateGraph 14노드 Chat Agent\n\n"
            "START → context_loader → route_has_image\n"
            "  이미지 → image_analyzer → intent_emotion_classifier\n"
            "  텍스트 → intent_emotion_classifier\n"
            "→ route_after_intent\n"
            "  recommend/search → preference_refiner\n"
            "    선호 부족 → question_generator → response_formatter\n"
            "    선호 충분 → query_builder → rag_retriever\n"
            "      품질 OK → llm_reranker → recommendation_ranker\n"
            "               → explanation_generator → response_formatter\n"
            "      품질 미달 → question_generator → response_formatter\n"
            "  general → general_responder → response_formatter\n"
            "  info/theater/booking → tool_executor_node → response_formatter"
        )

    # ── 슬라이드 9: 결과 제시 ③ ──────────────────────────────────────────────
    slide9 = prs.slides[8]
    s = find_shape_by_id(slide9, 16)  # TextBox 15
    if s:
        replace_text_direct(s, "* 결과 제시 ③ 추천 알고리즘")
    s = find_shape_by_id(slide9, 22)  # TextBox 1
    if s:
        replace_text_direct(s,
            "하이브리드 검색: Qdrant(벡터) + ES(BM25) + Neo4j(그래프) → RRF 합산 (k=60)\n\n"
            "CF+CBF 동적 가중치:\n"
            "  Cold Start → CBF 100%\n"
            "  Warm → CF 50% + CBF 50%\n"
            "  정상 → CF 60% + CBF 40%\n\n"
            "MMR 다양성 최적화: λ=0.7 (점수 0.7 + 다양성 0.3)\n"
            "Candidate 10~15편 → 최종 3~5편 선택\n\n"
            "Movie Match Agent: 두 영화 교집합 특성 기반 추천\n"
            "스코어: min(sim(candidate, A), sim(candidate, B))"
        )

    # ── 슬라이드 10: 결과 제시 ④ ─────────────────────────────────────────────
    slide10 = prs.slides[9]
    s = find_shape_by_id(slide10, 16)  # TextBox 15
    if s:
        replace_text_direct(s, "* 결과 제시 ④ 서비스 구현 현황")
    s = find_shape_by_id(slide10, 26)  # TextBox 20
    if s:
        replace_text_direct(s,
            "Phase 0~6 완료: 스캐폴딩→파이프라인→RAG→LLM체인→Chat Agent→추천엔진→Tools\n\n"
            "다국어 검색 ML-1~3: ES 영문 필드 + 한국어 키워드 200개 + 멀티턴 강화\n"
            "유저 활동 Phase 0~5: Like/EventLog/WatchHistory/BehaviorProfile\n"
            "포인트/결제/리워드 v3.2: 6등급 팝콘 테마, 1P=10원 통일\n"
            "Client: 250+ 모듈, 다크/라이트 모드, 반응형 3단계\n"
            "Admin: 749 모듈, 10탭, 42개 API\n"
            "Backend: 58개 테이블, 도메인별 API 완성\n\n"
            "311개 테스트 통과"
        )

    # ── 슬라이드 11: 결과 제시 ⑤ ─────────────────────────────────────────────
    slide11 = prs.slides[10]
    s = find_shape_by_id(slide11, 16)  # TextBox 15
    if s:
        replace_text_direct(s, "* 결과 제시 ⑤ 시연 및 데모")
    s = find_shape_by_id(slide11, 55)  # TextBox 54
    if s:
        replace_text_direct(s,
            "* 시연 항목:\n"
            "  1. AI 챗봇 대화 기반 영화 추천\n"
            "  2. 이미지 업로드 분석 추천\n"
            "  3. Movie Match (두 영화 교집합 추천)\n"
            "  4. 관리자 대시보드 (10탭)\n"
            "  5. 포인트/결제/등급 시스템"
        )

    # ── 슬라이드 12: 자체 평가 ────────────────────────────────────────────────
    slide12 = prs.slides[11]

    eval_replacements = {
        "사전 기획의 관점에서  프로젝트 결과물에 대한 완성도 평가(10점 만점)":
            "완성도: 7/10\n핵심 기능 모두 구현 완료",
        "프로젝트 결과물의  추후 개선점이나 보완할 점 등 내용 정리":
            "LoRA 파인튜닝, 운영 DB 재적재\n실시간 Toss Payments 연동",
        "프로젝트를 수행하면서\x0b느낀 점이나 경험한 성과(경력 계획 등과 연관)":
            "5개 서비스 통합 경험\n실제 LLM/벡터DB 운영 경험",
        "프로젝트를 수행하면서 느낀 점이나 경험한 성과(경력 계획 등과 연관)":
            "5개 서비스 통합 경험\n실제 LLM/벡터DB 운영 경험",
        "개인 또는 우리 팀이 잘한 부분과 아쉬운 점":
            "하이브리드 RAG 검색 구현\n포인트 경제 시스템 설계",
        "모델 평가 결과, 정확도가 00.00%로\x0b정확도 향상을 위해 모델 추후 개선 필요":
            "Phase ML-1~3 검색 품질 개선\n311개 테스트 통과",
        "모델 평가 결과, 정확도가 00.00%로 정확도 향상을 위해 모델 추후 개선 필요":
            "Phase ML-1~3 검색 품질 개선\n311개 테스트 통과",
    }
    for old, new in eval_replacements.items():
        replace_text_in_shape_on_slide(slide12, old, new)

    # 저장
    prs.save(str(TEMPLATE_DST))
    print(f"[작업 1] 완료: {TEMPLATE_DST}")


def replace_text_in_shape_on_slide(slide, old_text: str, new_text: str):
    """슬라이드 전체 shape에서 (그룹 포함) 텍스트 교체."""
    for shape in slide.shapes:
        replace_text_in_shape(shape, old_text, new_text)


# ─── 작업 2: 신규 버전 ───────────────────────────────────────────────────────

def rgb_to_hex(color) -> str:
    """
    RGBColor (bytes 서브클래스) 또는 (r,g,b) 튜플을 6자리 hex 문자열로 변환한다.
    RGBColor는 bytes이므로 인덱스([0][1][2])로 각 채널에 접근한다.
    """
    if isinstance(color, (list, tuple)):
        r, g, b = color[0], color[1], color[2]
    else:
        # RGBColor는 bytes 서브클래스
        r, g, b = color[0], color[1], color[2]
    return f"{r:02x}{g:02x}{b:02x}"


# 색상 팔레트
C_BG = RGBColor(0x1a, 0x1a, 0x2e)          # 다크 네이비 배경
C_HEADER = RGBColor(0xe9, 0x45, 0x60)      # 레드 강조
C_DEEP_BLUE = RGBColor(0x0f, 0x34, 0x60)   # 딥 블루
C_PURPLE = RGBColor(0x53, 0x34, 0x83)      # 퍼플
C_WHITE = RGBColor(0xff, 0xff, 0xff)       # 흰색
C_LIGHT_GRAY = RGBColor(0xe0, 0xe0, 0xe0) # 연회색
C_ACCENT = RGBColor(0x16, 0xc7, 0x9a)      # 민트 강조
C_DARK = RGBColor(0x16, 0x21, 0x3e)        # 더 어두운 배경

FONT_KR = "Malgun Gothic"
FONT_EN = "Calibri"

W = Inches(13.33)
H = Inches(7.5)


def add_background(slide, color: RGBColor):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, color: RGBColor = None):
    """상단 헤더 바 추가 (높이 0.5인치)."""
    if color is None:
        color = C_HEADER
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, font_name=None, wrap=True):
    """텍스트박스를 추가하고 텍스트/서식을 설정한다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # 기존 paragraph 정리
    from lxml import etree
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    lines = text.split("\n")
    for line in lines:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p_elem, qn("a:pPr"))
        align_map = {PP_ALIGN.CENTER: "ctr", PP_ALIGN.LEFT: "l", PP_ALIGN.RIGHT: "r"}
        pPr.set("algn", align_map.get(align, "l"))

        r_elem = etree.SubElement(p_elem, qn("a:r"))
        rpr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "ko-KR", "dirty": "0"})
        if bold:
            rpr.set("b", "1")
        rpr.set("sz", str(int(font_size * 100)))

        fn = font_name or FONT_KR
        lat = etree.SubElement(rpr, qn("a:latin"))
        lat.set("typeface", fn)

        if color:
            sf = etree.SubElement(rpr, qn("a:solidFill"))
            srg = etree.SubElement(sf, qn("a:srgbClr"))
            srg.set("val", rgb_to_hex(color))

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line

    return txBox


def add_colored_box(slide, left, top, width, height, fill_color, line_color=None):
    """채워진 사각형 도형을 추가한다."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_section_title(slide, title_text, top=Inches(0.6)):
    """섹션 제목 텍스트박스를 추가한다."""
    add_textbox(
        slide,
        left=Inches(0.5), top=top,
        width=Inches(12), height=Inches(0.6),
        text=title_text,
        font_size=28, bold=True,
        color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )
    # 제목 하단 구분선
    line = slide.shapes.add_shape(1, Inches(0.5), top + Inches(0.65), Inches(12), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_HEADER
    line.line.fill.background()


def add_slide_number(slide, num: int):
    """오른쪽 하단 슬라이드 번호."""
    add_textbox(
        slide,
        left=Inches(12.3), top=Inches(7.1),
        width=Inches(0.8), height=Inches(0.3),
        text=str(num),
        font_size=10, bold=False,
        color=C_LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )


def add_table_new(slide, left, top, width, height, headers, rows,
                  header_bg=None, row_bg=None, font_size=11):
    """
    신규 버전용 표 추가.
    headers: 헤더 텍스트 리스트
    rows: 데이터 행 리스트 (각 행은 열값 리스트)
    """
    n_cols = len(headers)
    n_rows = len(rows) + 1  # 헤더 포함

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    if header_bg is None:
        header_bg = C_DEEP_BLUE
    if row_bg is None:
        row_bg = RGBColor(0x0d, 0x1b, 0x35)

    # 헤더
    for col_idx, hdr in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        set_cell_text(cell, hdr, bold=True, font_size=font_size, color=(0xff, 0xff, 0xff))
        # 배경
        from lxml import etree
        tc = cell._tc
        tcPr = tc.find(qn("a:tcPr"))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn("a:tcPr"))
        # 배경색 설정
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", f"{header_bg.red:02x}{header_bg.green:02x}{header_bg.blue:02x}")

    # 데이터 행
    for row_idx, row_data in enumerate(rows):
        bg = row_bg if row_idx % 2 == 0 else RGBColor(0x12, 0x1e, 0x38)
        for col_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(row_idx + 1, col_idx)
            set_cell_text(cell, str(cell_text), font_size=font_size - 1,
                          color=(0xe0, 0xe0, 0xe0))
            from lxml import etree
            tc = cell._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn("a:tcPr"))
            solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
            srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgb.set("val", f"{bg.red:02x}{bg.green:02x}{bg.blue:02x}")

    return tbl_shape


def make_new_version():
    """
    완전히 새로운 디자인의 몽글픽 중간보고 PPTX를 생성한다.
    총 14슬라이드 구성.
    """
    print("[작업 2] 신규 버전 생성 시작...")

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 빈 레이아웃 사용 (인덱스 6 = blank)
    blank_layout = prs.slide_layouts[6]

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)

    # 중앙 상단 로고 영역 박스
    add_colored_box(slide, 0, 0, W, Inches(1.2), C_DEEP_BLUE)

    # 서비스명
    add_textbox(slide,
        left=Inches(0.5), top=Inches(0.2),
        width=Inches(8), height=Inches(0.8),
        text="MongLePick | 몽글픽",
        font_size=22, bold=True,
        color=C_HEADER,
        align=PP_ALIGN.LEFT,
    )

    # 날짜
    add_textbox(slide,
        left=Inches(10.5), top=Inches(0.3),
        width=Inches(2.5), height=Inches(0.5),
        text="2026. 04. 07",
        font_size=14, bold=False,
        color=C_LIGHT_GRAY,
        align=PP_ALIGN.RIGHT,
    )

    # 메인 제목
    add_textbox(slide,
        left=Inches(1), top=Inches(2.0),
        width=Inches(11), height=Inches(1.5),
        text="몽글픽 (MongLePick)",
        font_size=52, bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )

    # 부제목
    add_textbox(slide,
        left=Inches(1), top=Inches(3.5),
        width=Inches(11), height=Inches(0.8),
        text="AI 기반 개인화 영화 추천 서비스",
        font_size=28, bold=False,
        color=C_HEADER,
        align=PP_ALIGN.CENTER,
    )

    # 구분선
    line = slide.shapes.add_shape(1, Inches(3), Inches(4.4), Inches(7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = C_HEADER
    line.line.fill.background()

    # 팀 정보
    add_textbox(slide,
        left=Inches(1), top=Inches(4.6),
        width=Inches(11), height=Inches(0.5),
        text="팀 몽글 | K-Digital Training",
        font_size=18, bold=False,
        color=C_LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )

    # 팀원
    add_textbox(slide,
        left=Inches(1), top=Inches(5.2),
        width=Inches(11), height=Inches(0.5),
        text="윤형주 (팀장)  |  김민규  |  이민수  |  정한나",
        font_size=15, bold=False,
        color=C_LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )

    # 중간보고 배지
    badge = add_colored_box(slide, Inches(5.5), Inches(6.0), Inches(2.3), Inches(0.65), C_PURPLE)
    add_textbox(slide,
        left=Inches(5.5), top=Inches(6.05),
        width=Inches(2.3), height=Inches(0.55),
        text="중간 보고서",
        font_size=16, bold=True,
        color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_slide_number(slide, 1)

    # ── 슬라이드 2: 목차 ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "목차  Contents")
    add_slide_number(slide, 2)

    toc_items = [
        ("01", "프로젝트 개요", "서비스 소개, 선정 배경, 차별화 포인트"),
        ("02", "시스템 아키텍처", "5개 서비스 구조 + 인프라"),
        ("03", "AI Agent 설계", "14노드 LangGraph 흐름"),
        ("04", "데이터 파이프라인", "910,140건 5DB 적재"),
        ("05", "추천 알고리즘", "하이브리드 RAG + CF/CBF + MMR"),
        ("06", "팀 구성 및 역할", "4명 담당 영역"),
        ("07", "WBS / 개발 일정", "6단계 Phase 기반"),
        ("08", "구현 현황", "Phase 0~6 완료 사항"),
        ("09", "주요 기능 화면", "핵심 기능 목록"),
        ("10", "향후 계획 및 자체 평가", "남은 작업 + 완성도 평가"),
    ]

    col_left = [Inches(0.5), Inches(6.9)]
    for idx, (num, title, desc) in enumerate(toc_items):
        col = idx // 5
        row = idx % 5
        left = col_left[col]
        top = Inches(1.5) + Inches(1.0) * row

        # 번호 박스
        num_box = add_colored_box(slide, left, top, Inches(0.5), Inches(0.5),
                                   C_HEADER if idx % 2 == 0 else C_PURPLE)
        add_textbox(slide, left, top + Inches(0.05), Inches(0.5), Inches(0.4),
                    num, font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 제목
        add_textbox(slide,
            left=left + Inches(0.6), top=top,
            width=Inches(5.5), height=Inches(0.3),
            text=title,
            font_size=16, bold=True, color=C_WHITE)

        # 설명
        add_textbox(slide,
            left=left + Inches(0.6), top=top + Inches(0.28),
            width=Inches(5.5), height=Inches(0.28),
            text=desc,
            font_size=11, bold=False, color=C_LIGHT_GRAY)

    # ── 슬라이드 3: 프로젝트 개요 ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "01  프로젝트 개요")
    add_slide_number(slide, 3)

    # 3개 카드
    cards = [
        (C_DEEP_BLUE, "서비스 소개",
         "영화 추천에서 출발한\nAI 챗봇 서비스\n\n사용자와 자연어 대화로\n취향을 파악하고 최적의\n영화를 추천합니다.\n\n이미지 분석, 무드 기반 검색\n멀티턴 대화 지원"),
        (C_PURPLE, "선정 배경",
         "기존 추천 시스템의 한계:\n- 단순 평점 기반\n- 획일적 콘텐츠 필터링\n- 컨텍스트 미반영\n\n몽글픽의 접근:\n대화 기반 초개인화\nLLM + 하이브리드 RAG"),
        (C_DEEP_BLUE, "차별화 포인트",
         "하이브리드 RAG 검색\nQdrant+ES+Neo4j RRF\n\n14노드 LangGraph 에이전트\n이미지 분석 추천\nMovie Match 기능\n\n포인트/등급 경제 시스템\n6등급 팝콘 테마"),
    ]

    for i, (bg, title, body) in enumerate(cards):
        left = Inches(0.4 + i * 4.3)
        top = Inches(1.4)
        w = Inches(4.1)
        h = Inches(5.7)
        add_colored_box(slide, left, top, w, h, bg)
        # 카드 상단 강조 바
        add_colored_box(slide, left, top, w, Inches(0.1), C_HEADER)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                    w - Inches(0.4), Inches(0.5),
                    title, font_size=18, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.9),
                    w - Inches(0.4), h - Inches(1.2),
                    body, font_size=13, color=C_LIGHT_GRAY)

    # ── 슬라이드 4: 시스템 아키텍처 ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "02  시스템 아키텍처")
    add_slide_number(slide, 4)

    # 5개 서비스 박스
    services = [
        (C_HEADER, "AI Agent", ":8000", "FastAPI + LangGraph\n14노드 Chat Agent\nOllama LLM"),
        (C_DEEP_BLUE, "Backend", ":8080", "Spring Boot 3\nJPA + JWT/OAuth2\n58개 테이블"),
        (C_PURPLE, "Recommend", ":8001", "FastAPI\nCF/CBF 하이브리드\nSQLAlchemy + Redis"),
        (RGBColor(0x1a, 0x5c, 0x3a), "Client", ":5173", "React + Vite\nstyled-components\n250+ 모듈"),
        (RGBColor(0x5c, 0x3a, 0x1a), "Admin", ":5174", "React + Vite\n10탭 대시보드\n749 모듈"),
    ]

    for i, (color, name, port, desc) in enumerate(services):
        left = Inches(0.3 + i * 2.6)
        top = Inches(1.5)
        w = Inches(2.4)
        h = Inches(2.4)
        add_colored_box(slide, left, top, w, h, color)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.1),
                    w - Inches(0.2), Inches(0.45),
                    name, font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.5),
                    w - Inches(0.2), Inches(0.3),
                    port, font_size=13, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.9),
                    w - Inches(0.2), Inches(1.4),
                    desc, font_size=12, color=C_WHITE)

    # DB/인프라 섹션
    add_textbox(slide,
        left=Inches(0.3), top=Inches(4.1),
        width=Inches(12.5), height=Inches(0.35),
        text="데이터베이스 & 인프라",
        font_size=15, bold=True, color=C_HEADER)

    db_items = [
        (C_DEEP_BLUE, "MySQL 8.0\n58개 테이블"),
        (C_PURPLE, "Qdrant\n4096차원 벡터"),
        (RGBColor(0x8c, 0x3a, 0x00), "Elasticsearch 8.17\nNori 한국어"),
        (RGBColor(0x1a, 0x5c, 0x3c), "Neo4j 5\n그래프 DB"),
        (RGBColor(0x5c, 0x1a, 0x1a), "Redis 7\n세션 + 캐시"),
        (RGBColor(0x2a, 0x2a, 0x5c), "카카오 클라우드\n4-VM 인프라"),
    ]

    for i, (color, label) in enumerate(db_items):
        left = Inches(0.3 + i * 2.15)
        top = Inches(4.55)
        w = Inches(2.0)
        h = Inches(1.2)
        add_colored_box(slide, left, top, w, h, color)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.2),
                    w - Inches(0.2), h - Inches(0.3),
                    label, font_size=12, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 인프라 설명
    add_textbox(slide,
        left=Inches(0.3), top=Inches(5.9),
        width=Inches(12.5), height=Inches(0.5),
        text="카카오 클라우드 4-VM: VM1(Public, Nginx+React) / VM2(Spring Boot + FastAPI) / VM3(Prometheus+Grafana) / VM4-GPU(vLLM+DB)",
        font_size=11, color=C_LIGHT_GRAY)

    # ── 슬라이드 5: AI Agent 설계 ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "03  AI Agent 설계 (LangGraph 14노드)")
    add_slide_number(slide, 5)

    # LLM 정보 박스
    llm_info = [
        (C_HEADER, "EXAONE 4.0 32B", "선호 추출, 대화 생성, 추천 이유\ntemperature < 0.6 | Ollama"),
        (C_DEEP_BLUE, "qwen3.5:35b-a3b", "의도+감정 분류, 이미지 분석\ntemperature 0.1 | Ollama"),
        (C_PURPLE, "Upstage Solar", "분류/추출/설명 API\n4096차원 임베딩"),
    ]

    for i, (color, model, desc) in enumerate(llm_info):
        left = Inches(0.3 + i * 4.3)
        top = Inches(1.4)
        add_colored_box(slide, left, top, Inches(4.0), Inches(1.1), color)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                    Inches(3.7), Inches(0.4),
                    model, font_size=15, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.5),
                    Inches(3.7), Inches(0.5),
                    desc, font_size=11, color=C_LIGHT_GRAY)

    # 그래프 흐름도
    add_textbox(slide,
        left=Inches(0.3), top=Inches(2.7),
        width=Inches(12.5), height=Inches(0.35),
        text="Chat Agent 그래프 흐름",
        font_size=15, bold=True, color=C_ACCENT)

    flow_text = (
        "START → context_loader → route_has_image\n"
        "         이미지 있음 → image_analyzer → intent_emotion_classifier\n"
        "         이미지 없음 → intent_emotion_classifier\n"
        "→ route_after_intent\n"
        "   recommend/search → preference_refiner → route_after_preference\n"
        "         선호 부족 → question_generator → response_formatter → END\n"
        "         선호 충분 → query_builder → rag_retriever → retrieval_quality_checker\n"
        "                      품질 OK → llm_reranker → recommendation_ranker\n"
        "                               → explanation_generator → response_formatter → END\n"
        "                      품질 미달 → question_generator → response_formatter → END\n"
        "   general → general_responder → response_formatter → END\n"
        "   info/theater/booking → tool_executor_node → response_formatter → END"
    )

    add_colored_box(slide, Inches(0.3), Inches(3.15), Inches(12.5), Inches(3.8), C_DARK)
    add_textbox(slide,
        left=Inches(0.5), top=Inches(3.2),
        width=Inches(12.2), height=Inches(3.7),
        text=flow_text,
        font_size=12, color=C_ACCENT,
        font_name=FONT_EN)

    # ── 슬라이드 6: 데이터 파이프라인 ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "04  데이터 파이프라인")
    add_slide_number(slide, 6)

    # 수집 소스 → 처리 → 적재 흐름
    add_textbox(slide,
        left=Inches(0.3), top=Inches(1.4),
        width=Inches(3.5), height=Inches(0.35),
        text="데이터 수집 소스",
        font_size=14, bold=True, color=C_HEADER)

    sources = [
        "TMDB API\n영화 메타데이터",
        "Kaggle Dataset\n영화 평점/리뷰",
        "KOBIS\n한국 박스오피스",
        "KMDb\n한국영화 DB",
    ]
    for i, src in enumerate(sources):
        left = Inches(0.3 + i * 3.1)
        add_colored_box(slide, left, Inches(1.85), Inches(2.8), Inches(1.0), C_DEEP_BLUE)
        add_textbox(slide, left + Inches(0.1), Inches(1.95),
                    Inches(2.6), Inches(0.85),
                    src, font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 화살표 텍스트
    add_textbox(slide,
        left=Inches(0.3), top=Inches(3.0),
        width=Inches(12.5), height=Inches(0.3),
        text="정제 / 정규화 / Upstage 임베딩 (4096차원) / Nori 형태소 분석",
        font_size=13, bold=True, color=C_HEADER, align=PP_ALIGN.CENTER)

    # 적재 현황 표
    add_textbox(slide,
        left=Inches(0.3), top=Inches(3.4),
        width=Inches(12.5), height=Inches(0.35),
        text="5DB 적재 현황  |  총 910,140건",
        font_size=14, bold=True, color=C_ACCENT)

    add_table_new(
        slide,
        left=Inches(0.3), top=Inches(3.85),
        width=Inches(12.5), height=Inches(2.5),
        headers=["DB", "역할", "색인/설정", "규모"],
        rows=[
            ["Qdrant", "벡터 유사도 검색", "4096차원, Cosine Similarity", "910K+ 벡터"],
            ["Elasticsearch 8.17", "BM25 키워드 검색", "Nori 한국어 + 영문 필드", "910K+ 문서"],
            ["Neo4j 5", "그래프 관계 탐색", "감독/배우/장르/무드 노드", "다중 관계"],
            ["MySQL 8.0", "메타데이터 저장", "58개 테이블, JPA+MyBatis", "사용자/이력/결제"],
            ["Redis 7", "세션 + CF 캐시", "TTL 30일, 협업필터링 캐시", "멀티턴 대화"],
        ],
        font_size=11,
    )

    # ── 슬라이드 7: 추천 알고리즘 ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "05  추천 알고리즘")
    add_slide_number(slide, 7)

    # 좌측: 하이브리드 검색
    add_colored_box(slide, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.7), C_DARK)
    add_textbox(slide, Inches(0.5), Inches(1.5),
                Inches(5.7), Inches(0.4),
                "하이브리드 RAG 검색", font_size=16, bold=True, color=C_HEADER)

    rag_text = (
        "Qdrant (벡터 검색)\n"
        "  - Upstage 4096차원 임베딩\n"
        "  - Cosine 유사도\n\n"
        "Elasticsearch (BM25)\n"
        "  - Nori 한국어 형태소 분석\n"
        "  - 영문 필드 + 대안 제목\n\n"
        "Neo4j (그래프 탐색)\n"
        "  - 감독/배우/장르 관계\n"
        "  - 멀티홉 쿼리\n\n"
        "RRF 합산 (k=60)\n"
        "  Candidate 10~15편 → 최종 3~5편"
    )
    add_textbox(slide, Inches(0.5), Inches(2.0),
                Inches(5.7), Inches(4.9),
                rag_text, font_size=13, color=C_LIGHT_GRAY)

    # 우측: CF+CBF + MMR
    add_colored_box(slide, Inches(6.7), Inches(1.4), Inches(6.1), Inches(2.6), C_DARK)
    add_textbox(slide, Inches(6.9), Inches(1.5),
                Inches(5.7), Inches(0.4),
                "CF + CBF 동적 가중치", font_size=16, bold=True, color=C_ACCENT)

    cf_text = (
        "Cold Start  → CBF 100%\n"
        "Warm        → CF 50% + CBF 50%\n"
        "정상        → CF 60% + CBF 40%\n\n"
        "UserBehaviorProfile (매일 03:00 배치)\n"
        "Shannon Entropy 기반 taste_consistency"
    )
    add_textbox(slide, Inches(6.9), Inches(2.0),
                Inches(5.7), Inches(1.9),
                cf_text, font_size=13, color=C_LIGHT_GRAY)

    add_colored_box(slide, Inches(6.7), Inches(4.15), Inches(6.1), Inches(2.95), C_DARK)
    add_textbox(slide, Inches(6.9), Inches(4.25),
                Inches(5.7), Inches(0.4),
                "Movie Match Agent", font_size=16, bold=True, color=C_PURPLE)

    match_text = (
        "두 영화 교집합 특성 → 함께 볼 영화 5편 추천\n\n"
        "스코어링:\n"
        "  min(sim(candidate, movieA), sim(candidate, movieB))\n\n"
        "유사도 가중치:\n"
        "  genre(0.35) + mood(0.25)\n"
        "  + keyword(0.15) + vector(0.25)\n\n"
        "MMR λ=0.7 (점수 0.7 + 다양성 0.3)"
    )
    add_textbox(slide, Inches(6.9), Inches(4.7),
                Inches(5.7), Inches(2.3),
                match_text, font_size=13, color=C_LIGHT_GRAY)

    # ── 슬라이드 8: 팀 구성 ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "06  팀 구성 및 역할")
    add_slide_number(slide, 8)

    team = [
        (C_HEADER, "윤형주", "팀장 / AI Agent",
         "AI Agent (LangGraph 14노드)\n"
         "추천 알고리즘 (CF+CBF+MMR)\n"
         "포인트/결제/리워드 시스템\n"
         "관리자 페이지 (42 API)\n"
         "Client UI/UX (108 모듈)\n"
         "Movie Match Agent"),
        (C_DEEP_BLUE, "김민규", "백엔드 개발",
         "인증 (JWT/OAuth2)\n"
         "사용자 도메인\n"
         "플레이리스트 API\n"
         "대시보드/사용자 관리\n"
         "Spring Boot JPA"),
        (C_PURPLE, "이민수", "백엔드 개발",
         "커뮤니티 (게시판/리뷰)\n"
         "콘텐츠 관리 (신고/혐오표현)\n"
         "MyBatis 기반 도메인\n"
         "Admin 콘텐츠 탭"),
        (RGBColor(0x1a, 0x5c, 0x3a), "정한나", "데이터/추천",
         "FastAPI 추천 서버\n"
         "통계/분석 API\n"
         "데이터 파이프라인\n"
         "Admin 통계 탭"),
    ]

    for i, (color, name, role, tasks) in enumerate(team):
        left = Inches(0.3 + i * 3.25)
        top = Inches(1.45)
        w = Inches(3.0)
        h = Inches(5.6)
        add_colored_box(slide, left, top, w, h, color)
        add_colored_box(slide, left, top, w, Inches(0.08), C_WHITE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.15),
                    w - Inches(0.3), Inches(0.5),
                    name, font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.65),
                    w - Inches(0.3), Inches(0.35),
                    role, font_size=13, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)
        # 구분선
        sep = slide.shapes.add_shape(1, left + Inches(0.3),
                                      top + Inches(1.1),
                                      w - Inches(0.6), Inches(0.04))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        sep.fill.fore_color.rgb = C_WHITE
        sep.line.fill.background()
        add_textbox(slide, left + Inches(0.15), top + Inches(1.2),
                    w - Inches(0.3), h - Inches(1.4),
                    tasks, font_size=12, color=C_WHITE)

    # ── 슬라이드 9: WBS / 개발 일정 ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "07  WBS / 개발 일정")
    add_slide_number(slide, 9)

    add_table_new(
        slide,
        left=Inches(0.3), top=Inches(1.45),
        width=Inches(12.5), height=Inches(4.8),
        headers=["단계", "기간", "주요 활동", "산출물", "상태"],
        rows=[
            ["요구사항 분석", "3/3 ~ 3/7", "프로젝트 기획, DB 설계, API 설계", "요구사항 명세서", "완료"],
            ["데이터 수집/적재", "3/10 ~ 3/14", "TMDB/KOBIS/KMDb 수집, 5DB 적재", "910,140건 데이터", "완료"],
            ["핵심 기능 구현", "3/17 ~ 3/28", "AI Agent, Backend API, Client UI", "Phase 0~5 완료", "완료"],
            ["통합/고도화", "3/31 ~ 4/4", "추천 품질, 다크모드, 관리자 페이지", "ML Phase 1~3", "완료"],
            ["중간 보고", "4/7", "중간 발표", "현재 발표", "진행중"],
            ["최종 마무리", "4/8 ~ 4/18", "LoRA 파인튜닝, DB 재적재, E2E 테스트", "최종 결과물", "예정"],
        ],
        font_size=11,
    )

    # Gantt 간이 표현
    add_textbox(slide,
        left=Inches(0.3), top=Inches(6.35),
        width=Inches(12.5), height=Inches(0.5),
        text="총 개발 기간: 2026년 3월 3일 ~ 4월 18일 (약 7주) | 전체 Phase: 0~6 완료 + Phase 7~8 예정",
        font_size=12, color=C_LIGHT_GRAY)

    # ── 슬라이드 10: 구현 현황 ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "08  구현 현황")
    add_slide_number(slide, 10)

    phases = [
        (C_HEADER,    "Phase 0~2", "완료", "스캐폴딩 + 데이터 파이프라인 + RAG 검색"),
        (C_DEEP_BLUE, "Phase 3~4", "완료", "LLM 체인 7개 + Chat Agent 14노드"),
        (C_PURPLE,    "Phase 5",   "완료", "추천 엔진 (CF+CBF+MMR) + 추천 품질 개선"),
        (C_HEADER,    "Phase 6",   "완료", "LangChain Tools 7개 (TMDB/지도/OTT/그래프)"),
        (C_DEEP_BLUE, "Phase ML-1~3", "완료", "다국어 검색 + 멀티턴 강화 + 키워드 매핑 200개"),
        (C_PURPLE,    "유저 데이터 Phase 0~5", "완료", "Like/EventLog/WatchHistory/BehaviorProfile"),
        (RGBColor(0x1a, 0x5c, 0x3a), "포인트/결제/리워드", "완료", "v3.2: 6등급, 1P=10원, 구독 상품"),
        (RGBColor(0x1a, 0x5c, 0x3a), "Client/Admin", "완료", "250+/749 모듈, 다크모드, 반응형"),
        (RGBColor(0x8c, 0x3a, 0x00), "Phase ML-4", "진행중", "운영 서버 Qdrant/ES 재적재"),
        (RGBColor(0x8c, 0x3a, 0x00), "LoRA 파인튜닝", "예정", "EXAONE 1.2B 몽글이 페르소나"),
        (RGBColor(0x5c, 0x1a, 0x1a), "Phase 7~8", "미착수", "분석+로드맵 에이전트, 통합 테스트"),
        (RGBColor(0x5c, 0x1a, 0x1a), "Toss Payments", "미착수", "실제 PG 연동 (현재 목 대체)"),
    ]

    for i, (color, phase, status, desc) in enumerate(phases):
        col = i // 6
        row = i % 6
        left = Inches(0.3 + col * 6.4)
        top = Inches(1.45 + row * 0.95)

        add_colored_box(slide, left, top, Inches(0.08), Inches(0.8), color)
        add_textbox(slide, left + Inches(0.2), top,
                    Inches(2.3), Inches(0.4),
                    phase, font_size=14, bold=True, color=C_WHITE)

        status_color = (C_ACCENT if status == "완료"
                        else C_HEADER if status == "진행중"
                        else C_LIGHT_GRAY)
        add_textbox(slide, left + Inches(2.5), top,
                    Inches(1.0), Inches(0.35),
                    f"[{status}]", font_size=12, bold=True, color=status_color)

        add_textbox(slide, left + Inches(0.2), top + Inches(0.38),
                    Inches(5.7), Inches(0.45),
                    desc, font_size=12, color=C_LIGHT_GRAY)

    # ── 슬라이드 11: 주요 기능 화면 ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "09  주요 기능")
    add_slide_number(slide, 11)

    features = [
        (C_HEADER, "AI 챗봇 추천",
         "자연어 대화 기반 영화 추천\n"
         "멀티턴 대화 + 이미지 분석\n"
         "8종 SSE 이벤트 스트리밍\n"
         "추천 이유 설명 생성"),
        (C_DEEP_BLUE, "Movie Match",
         "두 영화 교집합 추천\n"
         "장르/무드/키워드/벡터 스코어\n"
         "MMR 다양성 최적화\n"
         "SSE 스트리밍 지원"),
        (C_PURPLE, "포인트/등급",
         "6등급 팝콘 테마 (알갱이→몽아일체)\n"
         "AI 이용권 3-소스 모델\n"
         "포인트 상점 + 구독 플랜\n"
         "1P = 10원 통일"),
        (RGBColor(0x1a, 0x5c, 0x3a), "커뮤니티",
         "영화 리뷰 + 게시판\n"
         "추천 이력 + 플레이리스트\n"
         "업적/도장깨기/월드컵\n"
         "ReviewVote + 댓글"),
        (RGBColor(0x5c, 0x3a, 0x1a), "관리자",
         "10탭 관리 대시보드\n"
         "사용자/콘텐츠/결제/통계\n"
         "42개 Admin API\n"
         "다크모드 + 반응형"),
        (RGBColor(0x8c, 0x3a, 0x00), "검색/필터",
         "하이브리드 RAG 검색\n"
         "다국어 영화 검색 (ML-1~3)\n"
         "동적 필터 (장르/국가/연도)\n"
         "LLM 재랭킹"),
    ]

    for i, (color, title, body) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 4.3)
        top = Inches(1.45 + row * 2.9)
        w = Inches(4.1)
        h = Inches(2.7)
        add_colored_box(slide, left, top, w, h, color)
        add_colored_box(slide, left, top, w, Inches(0.08), C_WHITE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.15),
                    w - Inches(0.3), Inches(0.4),
                    title, font_size=16, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.65),
                    w - Inches(0.3), h - Inches(0.85),
                    body, font_size=13, color=C_LIGHT_GRAY)

    # ── 슬라이드 12: 향후 계획 ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "10  향후 계획")
    add_slide_number(slide, 12)

    plans = [
        (C_HEADER, "Phase ML-4\n운영 DB 재적재",
         "Qdrant 임베딩 재적재 (910K 벡터)\n"
         "ES 인덱스 재생성 (영문 필드 포함)\n"
         "TMDB 이미지 재수집\n"
         "예상 소요: 3~5일 (GPU VM 활용)"),
        (C_DEEP_BLUE, "몽글이 LoRA\n파인튜닝",
         "EXAONE 4.0 1.2B 기반\n"
         "몽글이 페르소나 파인튜닝\n"
         "친근한 말투 + 영화 전문성\n"
         "예상 소요: 4~5일 (Tesla T4)"),
        (C_PURPLE, "E2E 통합 테스트\n& Toss Payments",
         "5개 서비스 통합 시나리오 테스트\n"
         "Toss Payments 실제 SDK 연동\n"
         "cancelPayment 구현\n"
         "현재 311개 단위 테스트 통과"),
    ]

    for i, (color, title, body) in enumerate(plans):
        left = Inches(0.4 + i * 4.3)
        top = Inches(1.5)
        w = Inches(4.1)
        h = Inches(5.4)
        add_colored_box(slide, left, top, w, h, color)
        add_colored_box(slide, left, top, w, Inches(0.1), C_WHITE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                    w - Inches(0.4), Inches(0.7),
                    title, font_size=18, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.1),
                    w - Inches(0.4), h - Inches(1.3),
                    body, font_size=14, color=C_LIGHT_GRAY)

    # ── 슬라이드 13: 자체 평가 ────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)
    add_header_bar(slide)
    add_section_title(slide, "11  자체 평가")
    add_slide_number(slide, 13)

    # 완성도 바
    add_textbox(slide,
        left=Inches(0.3), top=Inches(1.45),
        width=Inches(12.5), height=Inches(0.4),
        text="완성도: 7 / 10",
        font_size=22, bold=True, color=C_WHITE)

    # 바 배경
    add_colored_box(slide, Inches(0.3), Inches(2.0), Inches(12.5), Inches(0.45),
                     RGBColor(0x2a, 0x2a, 0x4a))
    # 채워진 바 (70%)
    add_colored_box(slide, Inches(0.3), Inches(2.0), Inches(12.5 * 0.7), Inches(0.45), C_HEADER)
    add_textbox(slide, Inches(0.35), Inches(2.02),
                Inches(8.0), Inches(0.38),
                "70%  핵심 기능 모두 구현 완료",
                font_size=14, bold=True, color=C_WHITE)

    evals = [
        (C_DEEP_BLUE, "우수한 점",
         "5개 서비스 완전 통합 아키텍처 구현\n"
         "하이브리드 RAG 3종 DB 실제 연동\n"
         "실제 LLM (EXAONE/Solar) 운영\n"
         "포인트 경제 시스템 v3.2 설계+구현\n"
         "311개 테스트 통과"),
        (C_PURPLE, "아쉬운 점",
         "운영 환경 E2E 통합 테스트 부족\n"
         "LoRA 파인튜닝 미완성\n"
         "Toss Payments 실제 연동 미완\n"
         "영화 이미지 크롤링 미완\n"
         "Phase 7~8 미착수"),
        (RGBColor(0x1a, 0x5c, 0x3a), "개선 계획",
         "Phase ML-4: 운영 DB 재적재\n"
         "EXAONE 1.2B LoRA 파인튜닝\n"
         "Toss Payments cancelPayment 구현\n"
         "5서비스 통합 E2E 시나리오 테스트\n"
         "최종 발표일까지 마무리 목표"),
    ]

    for i, (color, title, body) in enumerate(evals):
        left = Inches(0.3 + i * 4.3)
        top = Inches(2.6)
        w = Inches(4.1)
        h = Inches(4.5)
        add_colored_box(slide, left, top, w, h, color)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                    w - Inches(0.4), Inches(0.45),
                    title, font_size=18, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.7),
                    w - Inches(0.4), h - Inches(0.9),
                    body, font_size=13, color=C_LIGHT_GRAY)

    # ── 슬라이드 14: Q&A ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_background(slide, C_BG)

    # 상단 바
    add_colored_box(slide, 0, 0, W, Inches(1.0), C_DEEP_BLUE)
    add_textbox(slide,
        left=Inches(0.5), top=Inches(0.2),
        width=Inches(8), height=Inches(0.6),
        text="MongLePick | 몽글픽",
        font_size=22, bold=True, color=C_HEADER)

    # Q&A 텍스트
    add_textbox(slide,
        left=Inches(1), top=Inches(1.8),
        width=Inches(11), height=Inches(2.0),
        text="Q & A",
        font_size=80, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        left=Inches(1), top=Inches(3.8),
        width=Inches(11), height=Inches(0.7),
        text="감사합니다",
        font_size=32, bold=False, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # 구분선
    line = slide.shapes.add_shape(1, Inches(4), Inches(4.6), Inches(5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = C_HEADER
    line.line.fill.background()

    # 팀 정보
    add_textbox(slide,
        left=Inches(1), top=Inches(4.8),
        width=Inches(11), height=Inches(0.5),
        text="팀 몽글  |  윤형주 · 김민규 · 이민수 · 정한나",
        font_size=16, color=C_LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        left=Inches(1), top=Inches(5.4),
        width=Inches(11), height=Inches(0.4),
        text="GitHub: monglepick/monglepick-agent  |  2026. 04. 07",
        font_size=13, color=RGBColor(0x80, 0x80, 0xa0), align=PP_ALIGN.CENTER)

    add_slide_number(slide, 14)

    # 저장
    prs.save(str(NEW_DST))
    print(f"[작업 2] 완료: {NEW_DST}")


# ─── 메인 ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    make_template_version()
    make_new_version()
    print("\n두 파일 모두 생성 완료!")
    print(f"  1. {TEMPLATE_DST}")
    print(f"  2. {NEW_DST}")
