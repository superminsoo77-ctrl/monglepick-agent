"""
몽글픽 중간보고 신규버전 PPTX 생성 스크립트
완전히 새로운 다크 테마 디자인 (14슬라이드)

실행:
  PYTHONPATH=src uv run python scripts/make_new_pptx.py
"""

from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─── 경로 ────────────────────────────────────────────────────────────────────
OUTPUT = Path("/Users/yoonhyungjoo/Documents/monglepick/docs/중간보고_몽글픽_신규버전.pptx")

# ─── 색상 팔레트 ─────────────────────────────────────────────────────────────
# RGBColor는 bytes 서브클래스이므로 [0],[1],[2] 인덱스로 각 채널에 접근
C_BG        = RGBColor(0x1a, 0x1a, 0x2e)  # 다크 네이비
C_HEADER    = RGBColor(0xe9, 0x45, 0x60)  # 레드 강조
C_DEEP_BLUE = RGBColor(0x0f, 0x34, 0x60)  # 딥 블루
C_PURPLE    = RGBColor(0x53, 0x34, 0x83)  # 퍼플
C_WHITE     = RGBColor(0xff, 0xff, 0xff)  # 흰색
C_LGRAY     = RGBColor(0xd0, 0xd0, 0xd0)  # 연회색
C_ACCENT    = RGBColor(0x16, 0xc7, 0x9a)  # 민트
C_DARK      = RGBColor(0x0d, 0x17, 0x35)  # 가장 어두운 배경
C_GREEN     = RGBColor(0x1a, 0x5c, 0x3a)  # 그린
C_ORANGE    = RGBColor(0x8c, 0x45, 0x00)  # 오렌지
C_BROWN     = RGBColor(0x5c, 0x30, 0x1a)  # 브라운
C_DARK_RED  = RGBColor(0x5c, 0x1a, 0x1a)  # 다크 레드
C_DARK_BLUE2= RGBColor(0x1e, 0x2a, 0x4a)  # 테이블 짝수 행
C_DARK_BLUE3= RGBColor(0x15, 0x1f, 0x38)  # 테이블 홀수 행

FONT_KR = "Malgun Gothic"
FONT_EN = "Calibri"

# 슬라이드 크기 (와이드 16:9)
W = Inches(13.33)
H = Inches(7.5)


# ─── 헬퍼: RGBColor → hex 문자열 ─────────────────────────────────────────────
def h(color) -> str:
    """RGBColor(bytes) 또는 (r,g,b) 시퀀스를 6자리 hex 문자열로 변환."""
    return f"{color[0]:02x}{color[1]:02x}{color[2]:02x}"


# ─── 기본 도형 추가 ───────────────────────────────────────────────────────────
def bg(slide, color):
    """슬라이드 배경색을 단색으로 설정."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, w, ht, color, line=False):
    """채워진 사각형 도형 추가. line=False이면 테두리 없음."""
    shp = slide.shapes.add_shape(1, left, top, w, ht)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = C_WHITE
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    return shp


# ─── 텍스트박스 추가 ──────────────────────────────────────────────────────────
def tb(slide, left, top, w, ht, lines,
       size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
       font=None, wrap=True):
    """
    텍스트박스를 추가한다.
    lines: str (\\n으로 줄바꿈) 또는 list[str]
    각 줄을 별도 paragraph로 구성하여 줄 간격을 유지한다.
    """
    if color is None:
        color = C_WHITE
    if font is None:
        font = FONT_KR

    txBox = slide.shapes.add_textbox(left, top, w, ht)
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # 기존 paragraph 제거
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # 줄 분리
    if isinstance(lines, str):
        parts = lines.split("\n")
    else:
        parts = lines

    align_map = {
        PP_ALIGN.CENTER: "ctr",
        PP_ALIGN.LEFT:   "l",
        PP_ALIGN.RIGHT:  "r",
    }
    algn = align_map.get(align, "l")

    for text in parts:
        p_el = etree.SubElement(txBody, qn("a:p"))
        pPr  = etree.SubElement(p_el, qn("a:pPr"))
        pPr.set("algn", algn)

        r_el = etree.SubElement(p_el, qn("a:r"))
        rPr  = etree.SubElement(r_el, qn("a:rPr"),
                                attrib={"lang": "ko-KR", "dirty": "0"})
        if bold:
            rPr.set("b", "1")
        rPr.set("sz", str(int(size * 100)))

        lat = etree.SubElement(rPr, qn("a:latin"))
        lat.set("typeface", font)

        sf  = etree.SubElement(rPr, qn("a:solidFill"))
        srg = etree.SubElement(sf,  qn("a:srgbClr"))
        srg.set("val", h(color))

        t_el = etree.SubElement(r_el, qn("a:t"))
        t_el.text = text

    return txBox


# ─── 표 추가 ──────────────────────────────────────────────────────────────────
def tbl(slide, left, top, w, ht, headers, rows,
        hdr_color=None, font_size=11):
    """
    헤더 + 데이터 행으로 구성된 표를 추가한다.
    hdr_color: 헤더 행 배경 RGBColor (기본 C_DEEP_BLUE)
    홀짝 행 색상을 번갈아 적용한다.
    """
    if hdr_color is None:
        hdr_color = C_DEEP_BLUE

    ncols = len(headers)
    nrows = len(rows) + 1  # 헤더 포함

    tbl_shp = slide.shapes.add_table(nrows, ncols, left, top, w, ht)
    t = tbl_shp.table

    def _set_cell(cell, text, fsize, bold_flag, fg, bg_color):
        """셀 텍스트와 배경색을 설정하는 내부 함수."""
        tc = cell._tc
        # 배경색
        tcPr = tc.find(qn("a:tcPr"))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn("a:tcPr"))
        sf  = etree.SubElement(tcPr, qn("a:solidFill"))
        srg = etree.SubElement(sf,   qn("a:srgbClr"))
        srg.set("val", h(bg_color))

        # 텍스트
        txBody = cell.text_frame._txBody
        for p in txBody.findall(qn("a:p")):
            txBody.remove(p)

        p_el = etree.SubElement(txBody, qn("a:p"))
        pPr  = etree.SubElement(p_el, qn("a:pPr"))
        pPr.set("algn", "l")

        r_el = etree.SubElement(p_el, qn("a:r"))
        rPr  = etree.SubElement(r_el, qn("a:rPr"),
                                attrib={"lang": "ko-KR", "dirty": "0"})
        if bold_flag:
            rPr.set("b", "1")
        rPr.set("sz", str(int(fsize * 100)))

        lat = etree.SubElement(rPr, qn("a:latin"))
        lat.set("typeface", FONT_KR)

        sf2  = etree.SubElement(rPr, qn("a:solidFill"))
        srg2 = etree.SubElement(sf2, qn("a:srgbClr"))
        srg2.set("val", h(fg))

        t_el = etree.SubElement(r_el, qn("a:t"))
        t_el.text = str(text)

    # 헤더 행
    for ci, hdr in enumerate(headers):
        _set_cell(t.cell(0, ci), hdr, font_size, True, C_WHITE, hdr_color)

    # 데이터 행 (홀짝 색상)
    for ri, row in enumerate(rows):
        row_bg = C_DARK_BLUE2 if ri % 2 == 0 else C_DARK_BLUE3
        for ci, val in enumerate(row):
            _set_cell(t.cell(ri + 1, ci), val, font_size - 1, False, C_LGRAY, row_bg)

    return tbl_shp


# ─── 섹션 헤더 ────────────────────────────────────────────────────────────────
def section_hdr(slide, title, num):
    """
    슬라이드 상단 섹션 헤더:
      - 빨간 바 (0.5인치 높이)
    위에 흰색 제목 텍스트
    """
    box(slide, 0, 0, W, Inches(0.55), C_HEADER)
    tb(slide, Inches(0.4), Inches(0.05),
       Inches(11), Inches(0.45),
       title, size=22, bold=True, color=C_WHITE)
    # 슬라이드 번호 (우상단)
    tb(slide, Inches(12.5), Inches(0.1),
       Inches(0.7), Inches(0.35),
       str(num), size=13, color=C_WHITE, align=PP_ALIGN.RIGHT)


# ─── 슬라이드 생성 함수들 ─────────────────────────────────────────────────────

def slide_01_cover(prs, layout):
    """슬라이드 1: 표지."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)

    # 상단 짙은 바
    box(sl, 0, 0, W, Inches(1.1), C_DEEP_BLUE)
    tb(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.7),
       "MongLePick | 몽글픽", size=24, bold=True, color=C_HEADER)
    tb(sl, Inches(10.5), Inches(0.35), Inches(2.6), Inches(0.4),
       "2026. 04. 07", size=13, color=C_LGRAY, align=PP_ALIGN.RIGHT)

    # 중앙 제목
    tb(sl, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.6),
       "몽글픽 (MongLePick)", size=56, bold=True, color=C_WHITE,
       align=PP_ALIGN.CENTER)

    # 부제목
    tb(sl, Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.8),
       "AI 기반 개인화 영화 추천 서비스",
       size=28, color=C_HEADER, align=PP_ALIGN.CENTER)

    # 구분선
    box(sl, Inches(3.0), Inches(4.3), Inches(7.33), Inches(0.06), C_HEADER)

    # 팀 정보
    tb(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5),
       "팀 몽글  |  K-Digital Training",
       size=17, color=C_LGRAY, align=PP_ALIGN.CENTER)
    tb(sl, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.45),
       "윤형주 (팀장)  ·  김민규  ·  이민수  ·  정한나",
       size=15, color=C_LGRAY, align=PP_ALIGN.CENTER)

    # 중간보고 배지
    box(sl, Inches(5.4), Inches(5.85), Inches(2.53), Inches(0.65), C_PURPLE)
    tb(sl, Inches(5.4), Inches(5.92), Inches(2.53), Inches(0.5),
       "중간 보고서", size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs, layout):
    """슬라이드 2: 목차."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "목차  Contents", 2)

    items = [
        ("01", "프로젝트 개요",       "서비스 소개 / 선정 배경 / 차별화 포인트"),
        ("02", "시스템 아키텍처",     "5개 서비스 구조 + 인프라 4-VM"),
        ("03", "AI Agent 설계",       "14노드 LangGraph 흐름 + LLM 구성"),
        ("04", "데이터 파이프라인",   "910,140건 5DB 수집·정제·임베딩·적재"),
        ("05", "추천 알고리즘",       "하이브리드 RAG + CF/CBF + MMR"),
        ("06", "팀 구성 및 역할",     "4명 담당 도메인"),
        ("07", "WBS / 개발 일정",     "6단계 Phase 기반 점진적 구현"),
        ("08", "구현 현황",           "Phase 0~6 완료 사항 + 진행/예정"),
        ("09", "주요 기능",           "AI 챗봇 / Movie Match / 포인트 / 관리자"),
        ("10", "향후 계획 및 자체 평가", "남은 작업 + 완성도 7/10"),
    ]

    colors = [C_HEADER, C_PURPLE]
    col_x  = [Inches(0.5), Inches(6.9)]

    for idx, (num, title, desc) in enumerate(items):
        col  = idx // 5
        row  = idx % 5
        lft  = col_x[col]
        tp   = Inches(0.8) + Inches(1.2) * row
        c    = colors[idx % 2]

        # 번호 박스
        box(sl, lft, tp + Inches(0.05), Inches(0.5), Inches(0.5), c)
        tb(sl, lft, tp + Inches(0.1), Inches(0.5), Inches(0.42),
           num, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 제목 + 설명
        tb(sl, lft + Inches(0.6), tp + Inches(0.0),
           Inches(5.8), Inches(0.42),
           title, size=16, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.6), tp + Inches(0.42),
           Inches(5.8), Inches(0.42),
           desc, size=11, color=C_LGRAY)


def slide_03_overview(prs, layout):
    """슬라이드 3: 프로젝트 개요."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "01  프로젝트 개요", 3)

    cards = [
        (C_DEEP_BLUE, "서비스 소개",
         ["영화 추천에서 출발한 AI 챗봇 서비스",
          "",
          "자연어 대화로 취향을 파악하고",
          "최적의 영화를 추천합니다.",
          "",
          "이미지 업로드 분석 추천",
          "멀티턴 대화 + 무드 기반 검색",
          "SSE 스트리밍 실시간 응답"]),
        (C_PURPLE, "선정 배경",
         ["기존 추천 시스템의 한계",
          "  - 단순 평점·장르 기반 필터링",
          "  - 컨텍스트·감정 미반영",
          "  - 일방향 추천 (대화 없음)",
          "",
          "몽글픽의 접근",
          "  대화 기반 초개인화",
          "  LLM + 하이브리드 RAG"]),
        (C_DEEP_BLUE, "차별화 포인트",
         ["하이브리드 RAG 3종 DB 연동",
          "  Qdrant + ES + Neo4j → RRF",
          "",
          "14노드 LangGraph 에이전트",
          "Movie Match 교집합 추천",
          "이미지 분석 → 추천",
          "",
          "포인트/등급 경제 시스템",
          "  6등급 팝콘 테마"]),
    ]

    for i, (color, title, body) in enumerate(cards):
        lft = Inches(0.35 + i * 4.35)
        tp  = Inches(0.7)
        w   = Inches(4.1)
        ht  = Inches(6.6)
        box(sl, lft, tp, w, ht, color)
        box(sl, lft, tp, w, Inches(0.1), C_HEADER)  # 상단 강조 바
        tb(sl, lft + Inches(0.2), tp + Inches(0.2),
           w - Inches(0.4), Inches(0.5),
           title, size=18, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.2), tp + Inches(0.85),
           w - Inches(0.4), ht - Inches(1.1),
           "\n".join(body), size=13, color=C_LGRAY)


def slide_04_arch(prs, layout):
    """슬라이드 4: 시스템 아키텍처."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "02  시스템 아키텍처", 4)

    # 5개 서비스 박스
    svcs = [
        (C_HEADER,    "AI Agent",   ":8000",
         "FastAPI + LangGraph\n14노드 Chat Agent\nOllama LLM (EXAONE/qwen3.5)\nMovie Match Agent"),
        (C_DEEP_BLUE, "Backend",    ":8080",
         "Spring Boot 3\nJPA + JWT/OAuth2\n58개 테이블\n도메인별 API"),
        (C_PURPLE,    "Recommend",  ":8001",
         "FastAPI\nCF/CBF 하이브리드\nSQLAlchemy + Redis\n정한나 담당"),
        (C_GREEN,     "Client",     ":5173",
         "React + Vite\nstyled-components\n250+ 모듈\n다크/라이트 모드"),
        (C_BROWN,     "Admin",      ":5174",
         "React + Vite\n10탭 대시보드\n749 모듈\n42개 Admin API"),
    ]

    for i, (color, name, port, desc) in enumerate(svcs):
        lft = Inches(0.25 + i * 2.6)
        tp  = Inches(0.65)
        w   = Inches(2.45)
        ht  = Inches(2.8)
        box(sl, lft, tp, w, ht, color)
        tb(sl, lft + Inches(0.1), tp + Inches(0.12),
           w - Inches(0.2), Inches(0.45),
           name, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, lft + Inches(0.1), tp + Inches(0.55),
           w - Inches(0.2), Inches(0.3),
           port, size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)
        tb(sl, lft + Inches(0.1), tp + Inches(0.95),
           w - Inches(0.2), ht - Inches(1.1),
           desc, size=12, color=C_WHITE)

    # DB 섹션 제목
    tb(sl, Inches(0.25), Inches(3.6),
       Inches(12.8), Inches(0.38),
       "데이터베이스 & 인프라", size=15, bold=True, color=C_HEADER)

    dbs = [
        (C_DEEP_BLUE,              "MySQL 8.0\n58개 테이블\nJPA+MyBatis 하이브리드"),
        (C_PURPLE,                 "Qdrant\n4096차원 벡터\nCosine Similarity"),
        (RGBColor(0x8c, 0x3a, 0x00), "Elasticsearch 8.17\nNori 한국어\n다국어 검색"),
        (C_GREEN,                  "Neo4j 5\n그래프 DB\n감독/배우/장르 관계"),
        (RGBColor(0x5c, 0x1a, 0x1a), "Redis 7\n세션 + CF 캐시\nTTL 30일"),
        (RGBColor(0x2a, 0x2a, 0x5c), "카카오 클라우드\n4-VM 인프라\nGPU VM (T4)"),
    ]
    for i, (color, label) in enumerate(dbs):
        lft = Inches(0.25 + i * 2.18)
        tp  = Inches(4.05)
        w   = Inches(2.0)
        ht  = Inches(1.5)
        box(sl, lft, tp, w, ht, color)
        tb(sl, lft + Inches(0.1), tp + Inches(0.18),
           w - Inches(0.2), ht - Inches(0.25),
           label, size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 인프라 설명
    tb(sl, Inches(0.25), Inches(5.7),
       Inches(12.8), Inches(0.55),
       "VM1(Public-Nginx+React) / VM2(Spring Boot+FastAPI+Recommend) / "
       "VM3(Prometheus+Grafana) / VM4-GPU(vLLM+MySQL+Redis+Qdrant+Neo4j+ES)",
       size=11, color=C_LGRAY)


def slide_05_agent(prs, layout):
    """슬라이드 5: AI Agent 설계."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "03  AI Agent 설계 (LangGraph 14노드)", 5)

    # LLM 구성 3개 카드
    llms = [
        (C_HEADER,    "EXAONE 4.0 32B",
         "선호 추출 / 대화 생성 / 추천 이유\ntemperature < 0.6\nOllama (로컬 서빙)"),
        (C_DEEP_BLUE, "qwen3.5:35b-a3b",
         "의도+감정 분류 / 이미지 분석\ntemperature 0.1\nOllama (로컬 서빙)"),
        (C_PURPLE,    "Upstage Solar API",
         "분류 / 추출 / 설명 생성\n4096차원 임베딩\nhybrid API 모드"),
    ]
    for i, (color, model, desc) in enumerate(llms):
        lft = Inches(0.25 + i * 4.38)
        tp  = Inches(0.65)
        w   = Inches(4.15)
        ht  = Inches(1.2)
        box(sl, lft, tp, w, ht, color)
        tb(sl, lft + Inches(0.15), tp + Inches(0.1),
           w - Inches(0.3), Inches(0.42),
           model, size=15, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.15), tp + Inches(0.55),
           w - Inches(0.3), Inches(0.55),
           desc, size=11, color=C_LGRAY)

    # 그래프 흐름도 박스
    tb(sl, Inches(0.25), Inches(2.0),
       Inches(12.8), Inches(0.38),
       "Chat Agent 그래프 흐름 (15노드)", size=15, bold=True, color=C_ACCENT)

    box(sl, Inches(0.25), Inches(2.45), Inches(12.8), Inches(4.85), C_DARK)

    flow = [
        "START → context_loader → route_has_image",
        "         이미지 있음 → image_analyzer → intent_emotion_classifier",
        "         이미지 없음 → intent_emotion_classifier",
        "→ route_after_intent",
        "   [recommend/search]",
        "     → preference_refiner → route_after_preference",
        "         선호 부족  → question_generator → response_formatter → END",
        "         선호 충분  → query_builder → rag_retriever → retrieval_quality_checker",
        "                       품질 OK   → llm_reranker → recommendation_ranker",
        "                                 → explanation_generator → response_formatter → END",
        "                       품질 미달 → question_generator → response_formatter → END",
        "   [general]       → general_responder → response_formatter → END",
        "   [info/theater/booking] → tool_executor_node (7 Tools) → response_formatter → END",
    ]
    tb(sl, Inches(0.45), Inches(2.55),
       Inches(12.4), Inches(4.65),
       "\n".join(flow), size=12, color=C_ACCENT, font=FONT_EN)


def slide_06_data(prs, layout):
    """슬라이드 6: 데이터 파이프라인."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "04  데이터 파이프라인", 6)

    # 수집 소스 4개 카드
    tb(sl, Inches(0.25), Inches(0.65),
       Inches(4), Inches(0.38),
       "데이터 수집 소스", size=14, bold=True, color=C_HEADER)

    srcs = [
        ("TMDB API",       "영화 메타데이터\n포스터 / 줄거리 / 출연진"),
        ("Kaggle Dataset", "영화 평점 / 리뷰 데이터"),
        ("KOBIS",          "한국 박스오피스\n국내 개봉 정보"),
        ("KMDb",           "한국영화 DB\n한국어 상세 정보"),
    ]
    for i, (src_name, src_desc) in enumerate(srcs):
        lft = Inches(0.25 + i * 3.27)
        tp  = Inches(1.1)
        w   = Inches(3.05)
        ht  = Inches(1.3)
        box(sl, lft, tp, w, ht, C_DEEP_BLUE)
        tb(sl, lft + Inches(0.1), tp + Inches(0.1),
           w - Inches(0.2), Inches(0.4),
           src_name, size=14, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.1), tp + Inches(0.5),
           w - Inches(0.2), Inches(0.7),
           src_desc, size=11, color=C_LGRAY)

    # 화살표 텍스트
    box(sl, Inches(0.25), Inches(2.55), Inches(12.8), Inches(0.5), C_PURPLE)
    tb(sl, Inches(0.25), Inches(2.6), Inches(12.8), Inches(0.42),
       "정제 / 정규화  →  Upstage Solar 임베딩 (4096차원)  →  Nori 형태소 분석  →  5DB 적재",
       size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 5DB 적재 표
    tb(sl, Inches(0.25), Inches(3.15),
       Inches(5), Inches(0.38),
       "5DB 적재 현황  |  총 910,140건", size=14, bold=True, color=C_ACCENT)

    tbl(sl,
        Inches(0.25), Inches(3.6),
        Inches(12.8), Inches(3.5),
        headers=["DB", "역할", "색인 / 설정", "규모"],
        rows=[
            ["Qdrant",             "벡터 유사도 검색",
             "4096차원, Cosine Similarity, HNSW",    "910K+ 벡터"],
            ["Elasticsearch 8.17", "BM25 키워드 검색",
             "Nori 한국어 + 영문 필드 + 대안 제목",    "910K+ 문서"],
            ["Neo4j 5",            "그래프 관계 탐색",
             "감독/배우/장르/무드 노드, 멀티홉 쿼리",  "다중 관계"],
            ["MySQL 8.0",          "메타데이터 저장",
             "58개 테이블, JPA+MyBatis 하이브리드",    "사용자/결제/이력"],
            ["Redis 7",            "세션 + CF 캐시",
             "TTL 30일, 멀티턴 대화 세션",              "협업필터 캐시"],
        ],
        font_size=12,
    )


def slide_07_algo(prs, layout):
    """슬라이드 7: 추천 알고리즘."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "05  추천 알고리즘", 7)

    # 좌: 하이브리드 RAG
    box(sl, Inches(0.25), Inches(0.65), Inches(6.3), Inches(6.6), C_DARK)
    tb(sl, Inches(0.45), Inches(0.72),
       Inches(5.9), Inches(0.42),
       "하이브리드 RAG 검색", size=16, bold=True, color=C_HEADER)

    rag = [
        "Qdrant (벡터 검색)",
        "  · Upstage Solar 4096차원 임베딩",
        "  · Cosine 유사도 HNSW 인덱스",
        "",
        "Elasticsearch 8.17 (BM25)",
        "  · Nori 한국어 형태소 분석",
        "  · 영문 제목 + 다국어 대안 제목",
        "  · 키워드 한국어 매핑 200개 (ML-2)",
        "",
        "Neo4j 5 (그래프 탐색)",
        "  · 감독/배우/장르 관계 멀티홉",
        "  · 유사 감독 추천 경로",
        "",
        "RRF 합산 (k=60)",
        "  Candidate 10~15편 → 최종 3~5편",
        "  LLM 재랭킹 (llm_reranker 노드)",
    ]
    tb(sl, Inches(0.45), Inches(1.22),
       Inches(5.9), Inches(5.9),
       "\n".join(rag), size=13, color=C_LGRAY)

    # 우상: CF+CBF
    box(sl, Inches(6.75), Inches(0.65), Inches(6.3), Inches(3.05), C_DARK)
    tb(sl, Inches(6.95), Inches(0.72),
       Inches(5.9), Inches(0.42),
       "CF + CBF 동적 가중치", size=16, bold=True, color=C_ACCENT)

    cf = [
        "Cold Start  →  CBF 100%",
        "Warm        →  CF 50% + CBF 50%",
        "정상        →  CF 60% + CBF 40%",
        "",
        "UserBehaviorProfile (매일 03:00 배치)",
        "Shannon Entropy → taste_consistency",
        "RecommendationImpact 성과 추적",
    ]
    tb(sl, Inches(6.95), Inches(1.22),
       Inches(5.9), Inches(2.4),
       "\n".join(cf), size=13, color=C_LGRAY)

    # 우하: Movie Match
    box(sl, Inches(6.75), Inches(3.85), Inches(6.3), Inches(3.4), C_DARK)
    tb(sl, Inches(6.95), Inches(3.92),
       Inches(5.9), Inches(0.42),
       "Movie Match Agent", size=16, bold=True, color=C_PURPLE)

    mm = [
        "두 영화 교집합 특성 → 함께 볼 영화 5편 추천",
        "",
        "스코어링:",
        "  min(sim(candidate, movieA), sim(candidate, movieB))",
        "",
        "유사도 가중치:",
        "  genre(0.35) + mood(0.25)",
        "  + keyword(0.15) + vector(0.25)",
        "",
        "MMR λ=0.7  (점수 0.7 + 다양성 0.3)",
    ]
    tb(sl, Inches(6.95), Inches(4.45),
       Inches(5.9), Inches(2.7),
       "\n".join(mm), size=13, color=C_LGRAY)


def slide_08_team(prs, layout):
    """슬라이드 8: 팀 구성 및 역할."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "06  팀 구성 및 역할", 8)

    members = [
        (C_HEADER, "윤형주", "팀장 / AI Agent",
         ["AI Agent (LangGraph 14노드)",
          "추천 알고리즘 (CF+CBF+MMR)",
          "포인트/결제/리워드 시스템",
          "관리자 페이지 (42 API)",
          "Client UI/UX (108 모듈)",
          "Movie Match Agent"]),
        (C_DEEP_BLUE, "김민규", "백엔드 개발",
         ["인증 (JWT/OAuth2/소셜)",
          "사용자 도메인",
          "플레이리스트 API (9개)",
          "대시보드 / 사용자 관리",
          "Admin 대시보드 탭"]),
        (C_PURPLE, "이민수", "백엔드 개발",
         ["커뮤니티 (게시판/리뷰)",
          "콘텐츠 관리 (신고/혐오표현)",
          "MyBatis 기반 도메인",
          "Admin 콘텐츠 탭"]),
        (C_GREEN, "정한나", "데이터 / 추천",
         ["FastAPI 추천 서버",
          "통계 / 분석 API",
          "데이터 파이프라인",
          "Admin 통계 탭"]),
    ]

    for i, (color, name, role, tasks) in enumerate(members):
        lft = Inches(0.25 + i * 3.28)
        tp  = Inches(0.65)
        w   = Inches(3.05)
        ht  = Inches(6.6)
        box(sl, lft, tp, w, ht, color)
        box(sl, lft, tp, w, Inches(0.09), C_WHITE)  # 상단 흰 선
        tb(sl, lft + Inches(0.15), tp + Inches(0.15),
           w - Inches(0.3), Inches(0.55),
           name, size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(sl, lft + Inches(0.15), tp + Inches(0.72),
           w - Inches(0.3), Inches(0.38),
           role, size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)
        # 구분선
        box(sl, lft + Inches(0.25), tp + Inches(1.18),
            w - Inches(0.5), Inches(0.04), C_WHITE)
        tb(sl, lft + Inches(0.15), tp + Inches(1.3),
           w - Inches(0.3), ht - Inches(1.5),
           "\n".join(tasks), size=13, color=C_WHITE)


def slide_09_wbs(prs, layout):
    """슬라이드 9: WBS / 개발 일정."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "07  WBS / 개발 일정", 9)

    tbl(sl,
        Inches(0.25), Inches(0.7),
        Inches(12.8), Inches(5.5),
        headers=["단계", "기간", "주요 활동", "산출물", "상태"],
        rows=[
            ["요구사항 분석",   "3/3 ~ 3/7",
             "프로젝트 기획, DB 설계, API 설계",
             "요구사항 명세서", "완료"],
            ["데이터 수집/적재", "3/10 ~ 3/14",
             "TMDB/KOBIS/KMDb 수집, 5DB 적재",
             "910,140건 데이터", "완료"],
            ["핵심 기능 구현",  "3/17 ~ 3/28",
             "AI Agent, Backend API, Client UI",
             "Phase 0~5 완료",  "완료"],
            ["통합 / 고도화",   "3/31 ~ 4/4",
             "추천 품질, 다크모드, 관리자 페이지",
             "Phase ML-1~3",    "완료"],
            ["중간 보고",       "4/7",
             "중간 발표",
             "본 발표",         "진행중"],
            ["최종 마무리",     "4/8 ~ 4/18",
             "LoRA 파인튜닝, DB 재적재, E2E 테스트",
             "최종 결과물",     "예정"],
        ],
        font_size=12,
    )

    tb(sl, Inches(0.25), Inches(6.35),
       Inches(12.8), Inches(0.5),
       "총 개발 기간: 2026년 3월 3일 ~ 4월 18일 (약 7주)  "
       "|  전체 Phase: 0~6 완료  +  Phase 7~8 예정",
       size=12, color=C_LGRAY)


def slide_10_status(prs, layout):
    """슬라이드 10: 구현 현황."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "08  구현 현황", 10)

    phases = [
        # (색상, 이름, 상태, 설명)
        (C_ACCENT,    "Phase 0~2",              "완료", "스캐폴딩 + 데이터 파이프라인 + RAG 검색"),
        (C_ACCENT,    "Phase 3~4",              "완료", "LLM 체인 7개 + Chat Agent 14노드"),
        (C_ACCENT,    "Phase 5",                "완료", "추천 엔진 (CF+CBF+MMR) + 품질 개선"),
        (C_ACCENT,    "Phase 6",                "완료", "LangChain Tools 7개 (TMDB/지도/OTT/그래프)"),
        (C_ACCENT,    "Phase ML-1~3",           "완료", "다국어 검색 + 멀티턴 강화 + 키워드 200개"),
        (C_ACCENT,    "유저 데이터 Phase 0~5",  "완료", "Like/EventLog/WatchHistory/BehaviorProfile"),
        (C_ACCENT,    "포인트/결제/리워드",     "완료", "v3.2: 6등급 팝콘 테마, 1P=10원, 구독 상품"),
        (C_ACCENT,    "Client / Admin",          "완료", "250+/749 모듈, 다크모드, 반응형 3단계"),
        (C_HEADER,    "Phase ML-4",              "진행중", "운영 서버 Qdrant/ES 재적재 (예상 3~5일)"),
        (C_HEADER,    "LoRA 파인튜닝",           "진행중", "EXAONE 1.2B 몽글이 페르소나 (Tesla T4)"),
        (C_LGRAY,     "Phase 7~8",              "미착수", "분석+로드맵 에이전트, E2E 통합 테스트"),
        (C_LGRAY,     "Toss Payments",          "미착수", "실제 PG 연동 (현재 Mock 대체 중)"),
    ]

    for i, (color, phase, status, desc) in enumerate(phases):
        col = i // 6
        row = i % 6
        lft = Inches(0.25 + col * 6.5)
        tp  = Inches(0.68 + row * 1.08)

        # 왼쪽 컬러 바
        box(sl, lft, tp, Inches(0.09), Inches(0.88), color)

        tb(sl, lft + Inches(0.2), tp + Inches(0.02),
           Inches(3.8), Inches(0.42),
           phase, size=14, bold=True, color=C_WHITE)

        # 상태 텍스트
        st_color = (C_ACCENT  if status == "완료"
                    else C_HEADER if status == "진행중"
                    else C_LGRAY)
        tb(sl, lft + Inches(3.8), tp + Inches(0.02),
           Inches(1.3), Inches(0.38),
           f"[{status}]", size=12, bold=True, color=st_color)

        tb(sl, lft + Inches(0.2), tp + Inches(0.45),
           Inches(5.8), Inches(0.42),
           desc, size=12, color=C_LGRAY)


def slide_11_features(prs, layout):
    """슬라이드 11: 주요 기능."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "09  주요 기능", 11)

    feats = [
        (C_HEADER, "AI 챗봇 추천",
         ["자연어 대화 기반 영화 추천",
          "멀티턴 대화 + 이미지 분석",
          "8종 SSE 이벤트 스트리밍",
          "추천 이유 설명 생성",
          "7개 LangChain Tools 연동"]),
        (C_DEEP_BLUE, "Movie Match",
         ["두 영화 교집합 추천",
          "장르/무드/키워드/벡터 스코어",
          "MMR 다양성 최적화",
          "SSE 스트리밍 지원",
          "5편 최종 추천 결과"]),
        (C_PURPLE, "포인트 / 등급",
         ["6등급 팝콘 테마",
          "  알갱이→강냉이→팝콘",
          "  카라멜팝콘→몽글팝콘→몽아일체",
          "AI 이용권 3-소스 모델",
          "1P = 10원 통일 (v3.2)"]),
        (C_GREEN, "커뮤니티",
         ["영화 리뷰 + 게시판",
          "추천 이력 / 플레이리스트",
          "업적 / 도장깨기 / 월드컵",
          "ReviewVote (helpful/unhelpful)",
          "댓글 + 신고 시스템"]),
        (C_ORANGE, "관리자 대시보드",
         ["10탭 완전 구현",
          "사용자/콘텐츠/결제/통계",
          "42개 Admin API",
          "다크모드 + 반응형 3단계",
          "749 모듈 빌드 성공"]),
        (C_BROWN, "하이브리드 검색",
         ["Qdrant + ES + Neo4j → RRF",
          "다국어 영화 검색 (ML-1~3)",
          "동적 필터 (장르/국가/연도)",
          "LLM 재랭킹 노드",
          "311개 테스트 통과"]),
    ]

    for i, (color, title, body) in enumerate(feats):
        col = i % 3
        row = i // 3
        lft = Inches(0.25 + col * 4.38)
        tp  = Inches(0.65 + row * 3.3)
        w   = Inches(4.15)
        ht  = Inches(3.1)
        box(sl, lft, tp, w, ht, color)
        box(sl, lft, tp, w, Inches(0.09), C_WHITE)
        tb(sl, lft + Inches(0.15), tp + Inches(0.15),
           w - Inches(0.3), Inches(0.45),
           title, size=17, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.15), tp + Inches(0.7),
           w - Inches(0.3), ht - Inches(0.9),
           "\n".join(body), size=13, color=C_LGRAY)


def slide_12_plan(prs, layout):
    """슬라이드 12: 향후 계획."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "10  향후 계획", 12)

    plans = [
        (C_HEADER, "Phase ML-4\n운영 DB 재적재",
         ["Qdrant 임베딩 재적재 (910K 벡터)",
          "ES 인덱스 재생성 (영문 필드 포함)",
          "TMDB 영화 이미지 재수집",
          "",
          "예상 소요: 3~5일",
          "  (GPU VM4, Tesla T4 활용)",
          "",
          "완료 시 다국어 검색 ML-1~3",
          "코드 수정이 실제로 적용됨"]),
        (C_DEEP_BLUE, "몽글이 LoRA\n파인튜닝",
         ["EXAONE 4.0 1.2B 기반",
          "몽글이 페르소나 파인튜닝",
          "  친근한 말투 + 영화 전문성",
          "",
          "예상 소요: 4~5일",
          "  (Tesla T4, vLLM 서빙)",
          "",
          "현재: EXAONE 4.0 32B로",
          "  임시 대체 운영 중"]),
        (C_PURPLE, "E2E 통합 테스트\n& Toss Payments",
         ["5개 서비스 통합 시나리오 테스트",
          "Toss Payments 실제 SDK 연동",
          "  cancelPayment 구현",
          "",
          "현재 311개 단위 테스트 통과",
          "통합 E2E 시나리오 미완",
          "",
          "최종 발표일까지",
          "  마무리 목표"]),
    ]

    for i, (color, title, body) in enumerate(plans):
        lft = Inches(0.3 + i * 4.38)
        tp  = Inches(0.65)
        w   = Inches(4.15)
        ht  = Inches(6.6)
        box(sl, lft, tp, w, ht, color)
        box(sl, lft, tp, w, Inches(0.1), C_WHITE)
        tb(sl, lft + Inches(0.2), tp + Inches(0.18),
           w - Inches(0.4), Inches(0.75),
           title, size=18, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.2), tp + Inches(1.05),
           w - Inches(0.4), ht - Inches(1.25),
           "\n".join(body), size=13, color=C_LGRAY)


def slide_13_eval(prs, layout):
    """슬라이드 13: 자체 평가."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)
    section_hdr(sl, "11  자체 평가", 13)

    # 완성도 바
    tb(sl, Inches(0.25), Inches(0.65),
       Inches(12.8), Inches(0.45),
       "완성도: 7 / 10  (핵심 기능 모두 구현 완료, LoRA/운영DB 미완)",
       size=18, bold=True, color=C_WHITE)

    bar_w = Inches(12.8)
    box(sl, Inches(0.25), Inches(1.18), bar_w, Inches(0.42),
        RGBColor(0x2a, 0x2a, 0x4a))  # 배경
    box(sl, Inches(0.25), Inches(1.18), Inches(12.8 * 0.7), Inches(0.42), C_HEADER)
    tb(sl, Inches(0.35), Inches(1.2), Inches(8.5), Inches(0.38),
       "70%", size=14, bold=True, color=C_WHITE)

    evals = [
        (C_DEEP_BLUE, "우수한 점",
         ["5개 서비스 완전 통합 아키텍처",
          "하이브리드 RAG 3종 DB 실제 연동",
          "실제 LLM (EXAONE/Solar) 운영",
          "포인트 경제 시스템 v3.2 설계+구현",
          "311개 단위 테스트 통과"]),
        (C_PURPLE, "아쉬운 점",
         ["운영 환경 E2E 통합 테스트 부족",
          "LoRA 파인튜닝 미완성",
          "Toss Payments 실제 연동 미완",
          "영화 이미지 크롤링 미완",
          "Phase 7~8 미착수"]),
        (C_GREEN, "개선 계획",
         ["Phase ML-4: 운영 DB 재적재",
          "EXAONE 1.2B LoRA 파인튜닝",
          "Toss Payments cancelPayment",
          "5서비스 통합 E2E 시나리오",
          "최종 발표일까지 마무리"]),
    ]

    for i, (color, title, body) in enumerate(evals):
        lft = Inches(0.25 + i * 4.38)
        tp  = Inches(1.75)
        w   = Inches(4.15)
        ht  = Inches(5.5)
        box(sl, lft, tp, w, ht, color)
        tb(sl, lft + Inches(0.2), tp + Inches(0.15),
           w - Inches(0.4), Inches(0.45),
           title, size=18, bold=True, color=C_WHITE)
        tb(sl, lft + Inches(0.2), tp + Inches(0.7),
           w - Inches(0.4), ht - Inches(0.9),
           "\n".join(body), size=13, color=C_LGRAY)


def slide_14_qa(prs, layout):
    """슬라이드 14: Q&A."""
    sl = prs.slides.add_slide(layout)
    bg(sl, C_BG)

    # 상단 바
    box(sl, 0, 0, W, Inches(1.0), C_DEEP_BLUE)
    tb(sl, Inches(0.5), Inches(0.2), Inches(9), Inches(0.65),
       "MongLePick | 몽글픽", size=24, bold=True, color=C_HEADER)

    # Q&A 대형 텍스트
    tb(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2),
       "Q & A", size=80, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    tb(sl, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.75),
       "감사합니다", size=34, color=C_LGRAY, align=PP_ALIGN.CENTER)

    # 구분선
    box(sl, Inches(3.8), Inches(4.6), Inches(5.73), Inches(0.06), C_HEADER)

    # 팀 정보
    tb(sl, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.5),
       "팀 몽글  |  윤형주 · 김민규 · 이민수 · 정한나",
       size=17, color=C_LGRAY, align=PP_ALIGN.CENTER)

    tb(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4),
       "GitHub: monglepick/monglepick-agent  |  2026. 04. 07",
       size=13, color=RGBColor(0x70, 0x70, 0x90), align=PP_ALIGN.CENTER)

    # 슬라이드 번호
    tb(sl, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35),
       "14", size=13, color=C_LGRAY, align=PP_ALIGN.RIGHT)


# ─── 메인 ───────────────────────────────────────────────────────────────────

def main():
    print("[작업 2] 신규 버전 생성 시작...")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # 빈 레이아웃 (인덱스 6 = Blank)
    layout = prs.slide_layouts[6]

    slide_01_cover(prs, layout)
    slide_02_toc(prs, layout)
    slide_03_overview(prs, layout)
    slide_04_arch(prs, layout)
    slide_05_agent(prs, layout)
    slide_06_data(prs, layout)
    slide_07_algo(prs, layout)
    slide_08_team(prs, layout)
    slide_09_wbs(prs, layout)
    slide_10_status(prs, layout)
    slide_11_features(prs, layout)
    slide_12_plan(prs, layout)
    slide_13_eval(prs, layout)
    slide_14_qa(prs, layout)

    prs.save(str(OUTPUT))
    print(f"[작업 2] 완료: {OUTPUT}")
    print(f"  총 슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
